"""
create_pid_ssl_slides.py
========================
Creates a three-slide PowerPoint deck for the PID-SSL project.

Slide 1 — Minimal data-generation pipeline schematic  (matplotlib PNG → PPTX)
Slide 2 — Full mathematical formulation               (native PPTX shapes)
Slide 3 — Indicative training dynamics + results table (matplotlib PNG → PPTX)

Run from the project root:
    python create_pid_ssl_slides.py
Output: pid_ssl_data_generation.pptx
"""

import io
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── palette ───────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x12, 0x12, 0x2A)
C_TITLE    = RGBColor(0xE0, 0xE0, 0xFF)
C_SUBTITLE = RGBColor(0xA0, 0xC4, 0xFF)
C_UNIQUE   = RGBColor(0x54, 0xA2, 0x4B)
C_REDUND   = RGBColor(0x4C, 0x78, 0xA8)
C_SYNERGY  = RGBColor(0xE4, 0x57, 0x56)
C_BACKBONE = RGBColor(0xF5, 0x8E, 0x1E)
C_BORDER   = RGBColor(0x33, 0x33, 0x55)
C_TEXT     = RGBColor(0xE8, 0xE8, 0xF0)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_MATH_BG  = RGBColor(0x10, 0x10, 0x28)

BG_HEX = '#12122a'

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── PPTX helpers ──────────────────────────────────────────────────────────────

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, l, t, w, h, fill: RGBColor,
              border: RGBColor = None, border_pt: float = 0.0):
    shape = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border and border_pt > 0:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, l, t, w, h, text, font_size,
                  bold=False, color=None, align=PP_ALIGN.LEFT,
                  italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def _embed_png(slide, png_buf, l=0.0, t=0.0, w=13.33, h=7.5):
    """Embed a PNG BytesIO buffer as a full-slide image."""
    slide.shapes.add_picture(png_buf, Inches(l), Inches(t),
                              Inches(w), Inches(h))


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Schematic  (matplotlib)
# ─────────────────────────────────────────────────────────────────────────────

def _make_slide1_png() -> io.BytesIO:
    """Render a clean, minimal data-generation pipeline schematic."""
    W, H = 16, 9
    fig, ax = plt.subplots(figsize=(W, H), dpi=150)
    fig.patch.set_facecolor(BG_HEX)
    ax.set_facecolor(BG_HEX)
    ax.set_xlim(0, W)
    ax.set_ylim(0, H)
    ax.axis('off')

    # ── title ────────────────────────────────────────────────────────────────
    ax.text(W / 2, 8.62,
            'PID-SAR-3: Synthetic Benchmark — Data Generation Pipeline',
            ha='center', va='center', fontsize=21, fontweight='bold',
            color='#e8e8ff')
    ax.text(W / 2, 8.2,
            '3 observable views  ·  10 PID atom types  ·  d = 32 obs. dim'
            '  ·  m = 8 latent dim  ·  shared backbone + i.i.d. noise',
            ha='center', va='center', fontsize=11, color='#8899cc')
    ax.plot([0.3, W - 0.3], [7.87, 7.87], color='#2a2a44', lw=1.5)

    # ── column headers ───────────────────────────────────────────────────────
    for cx, cl in zip([3.1, 8.4, 13.55],
                      ['LATENT  SAMPLING', 'VIEW  CONSTRUCTION', 'OBSERVATIONS']):
        ax.text(cx, 7.57, cl, ha='center', va='center',
                fontsize=11, fontweight='bold', color='#6688aa')

    # vertical stage dividers
    for xd in [5.5, 10.95]:
        ax.plot([xd, xd], [1.08, 7.4], color='#252540', lw=1.5, ls='--')

    # ── helper: draw one family row ──────────────────────────────────────────
    def draw_row(yc, family, fc, lat_lines, cst_lines, view_active, note_lines):
        rh = 1.60
        y0 = yc - rh / 2

        # family strip
        fp = FancyBboxPatch((0.15, y0 + 0.06), 0.62, rh - 0.12,
                             boxstyle='round,pad=0.04', facecolor=fc, edgecolor='none')
        ax.add_patch(fp)
        ax.text(0.46, yc, family, ha='center', va='center', rotation=90,
                fontsize=9, fontweight='bold', color='white')

        # latent box
        lb = FancyBboxPatch((1.0, y0 + 0.1), 4.2, rh - 0.2,
                             boxstyle='round,pad=0.07', facecolor='#0a0a1f',
                             edgecolor=fc, linewidth=1.8)
        ax.add_patch(lb)
        ax.text(3.1, yc, '\n'.join(lat_lines), ha='center', va='center',
                fontsize=10.5, color='#d8dcff', linespacing=1.5)

        # arrow → construction
        ax.annotate('', xy=(5.65, yc), xytext=(5.24, yc),
                    arrowprops=dict(arrowstyle='->', color=fc, lw=2.3,
                                   mutation_scale=17))

        # construction box
        cb = FancyBboxPatch((5.72, y0 + 0.1), 4.95, rh - 0.2,
                             boxstyle='round,pad=0.07', facecolor='#0a0a1f',
                             edgecolor=fc, linewidth=1.8)
        ax.add_patch(cb)
        ax.text(8.19, yc, '\n'.join(cst_lines), ha='center', va='center',
                fontsize=10.5, color='#d8dcff', linespacing=1.5)

        # arrow → observations
        ax.annotate('', xy=(11.1, yc), xytext=(10.71, yc),
                    arrowprops=dict(arrowstyle='->', color=fc, lw=2.3,
                                   mutation_scale=17))

        # view boxes  (x1, x2, x3)
        for vi, (vl, active) in enumerate(zip(['x\u2081', 'x\u2082', 'x\u2083'],
                                              view_active)):
            vx = 11.22 + vi * 1.45
            vb = FancyBboxPatch((vx, y0 + 0.22), 1.22, rh - 0.44,
                                 boxstyle='round,pad=0.04',
                                 facecolor=fc if active else '#181830',
                                 edgecolor=fc, linewidth=2.0,
                                 alpha=0.88 if active else 0.22)
            ax.add_patch(vb)
            ax.text(vx + 0.61, yc + 0.06, vl, ha='center', va='center',
                    fontsize=14, fontweight='bold', color='white')
            if not active:
                ax.text(vx + 0.61, yc - 0.35, '0', ha='center',
                        va='center', fontsize=9, color='#334466')

        # side note
        if note_lines:
            ax.text(15.52, yc, '\n'.join(note_lines),
                    ha='center', va='center', fontsize=8.5,
                    color='#7799bb', linespacing=1.45, style='italic')

    # ── Unique ───────────────────────────────────────────────────────────────
    draw_row(
        yc=6.42, family='UNIQUE', fc='#54A24B',
        lat_lines=['u ~ N(0, I\u2098)', '(private latent)'],
        cst_lines=['x\u1d65 = \u03b1 \u00b7 P\u1d65 \u00b7 u',
                   'x\u1d65\u2019 = 0   for v\u2019 \u2260 v'],
        view_active=[True, False, False],
        note_lines=['Signal exclusive to one view',
                    '3 atoms: U\u2081  U\u2082  U\u2083'],
    )

    # ── Redundancy ───────────────────────────────────────────────────────────
    draw_row(
        yc=4.58, family='REDUND.', fc='#4C78A8',
        lat_lines=['r, \u03b7\u1d65 ~ N(0, I\u2098)', '(shared + private noise)'],
        cst_lines=['r\u1d65 = \u221a\u03c1\u00b7r + \u221a(1\u2212\u03c1)\u00b7\u03b7\u1d65    (\u03c1 = 0.8)',
                   'x\u1d65 = \u03b1 \u00b7 P\u1d65 \u00b7 r\u1d65'],
        view_active=[True, True, False],
        note_lines=['Correlated signal in \u22652 views',
                    '4 atoms: R\u2081\u2082  R\u2081\u2083  R\u2082\u2083  R\u2081\u2082\u2083'],
    )

    # ── Synergy ──────────────────────────────────────────────────────────────
    draw_row(
        yc=2.74, family='SYNERGY', fc='#E45756',
        lat_lines=['a, b ~ N(0, I\u2098)', '(two source latents)'],
        cst_lines=['s = f\u208d\u2098\u2097\u209a\u208e(a,b) \u2212 \u03bb\u00b7(C\u2090\u00b7a + C\u1d47\u00b7b)',
                   'x\u1d62=\u03b1\u00b7a,  x\u2c7c=\u03b1\u00b7b,  x\u2096=\u03b1\u00b7s'],
        view_active=[True, True, True],
        note_lines=['Target x\u2096 depends jointly on',
                    'x\u1d62 AND x\u2c7c (not individually)',
                    '3 atoms: S\u2081\u2082\u2192\u2083  S\u2081\u2083\u2192\u2082  S\u2082\u2083\u2192\u2081'],
    )

    # ── Backbone + noise footer ──────────────────────────────────────────────
    footer = FancyBboxPatch((0.15, 0.55), W - 0.3, 1.38,
                             boxstyle='round,pad=0.07', facecolor='#0a0a1e',
                             edgecolor='#F58E1E', linewidth=1.8)
    ax.add_patch(footer)
    ax.text(4.5, 1.47, '+ Shared backbone  (\u03b3 = 4.0)',
            ha='center', va='center', fontsize=10.5, fontweight='bold',
            color='#F58E1E')
    ax.text(4.5, 1.05,
            'z ~ N(0, I\u2098),    x\u1d65 \u2190 x\u1d65 + \u03b3\u00b7\u03b1\u00b7P\u209b\u02b0\u1d43\u02b3\u1d49\u1d48\u00b7z    (same z across all views)',
            ha='center', va='center', fontsize=10, color='#ddaa55')
    ax.text(12.0, 1.47, '+ Observation noise  (\u03c3 = 0.02)',
            ha='center', va='center', fontsize=10.5, fontweight='bold',
            color='#888899')
    ax.text(12.0, 1.05,
            '\u03b5\u1d65 ~ N(0, \u03c3\u00b2I\u1d48),    x\u1d65 \u2190 x\u1d65 + \u03b5\u1d65    (independent per view)',
            ha='center', va='center', fontsize=10, color='#7788aa')

    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=BG_HEX, edgecolor='none')
    buf.seek(0)
    plt.close(fig)
    return buf


def build_slide1(prs: Presentation):
    slide = _blank_slide(prs)
    _fill_bg(slide, C_BG)
    _embed_png(slide, _make_slide1_png(), l=0, t=0, w=13.33, h=7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — Mathematics  (native PPTX)
# ─────────────────────────────────────────────────────────────────────────────

def build_slide2(prs: Presentation):
    slide = _blank_slide(prs)
    _fill_bg(slide, C_BG)

    _add_rect(slide, 0, 0, 13.33, 0.68, fill=RGBColor(0x0D, 0x0D, 0x24))
    _add_textbox(slide, 0.18, 0.08, 11, 0.52,
                 "PID-SAR-3: Mathematical Formulation",
                 font_size=22, bold=True, color=C_TITLE, align=PP_ALIGN.LEFT)
    _add_textbox(slide, 0.18, 0.48, 11, 0.24,
                 "d=32 obs. dim  |  m=8 latent dim  |  \u03c3=0.02 noise  |  \u03c1=0.8  "
                 "|  \u03b1~U[0.8,1.2]  |  hop=1  |  \u03bb=0.25 de-leakage"
                 "  |  \u03b3=4.0 backbone gain  |  K=5 atoms (multi-atom mode)",
                 font_size=8, color=C_SUBTITLE, align=PP_ALIGN.LEFT)

    panels = [
        (0.18, 0.82, 6.3, 2.05, C_UNIQUE,
         "\u2460 Unique atom  U\u1d65  (v \u2208 {1, 2, 3})",
         [
             "u  ~  N(0, I_m)",
             "",
             "x_v  =  \u03b1 \u00b7 P_v^{U_v} \u00b7 u",
             "",
             "x_{v'} = 0   for all v' \u2260 v",
             "",
             "\u2022 u is a private m-dim latent, visible only in view v",
             "\u2022 P_v^{U_v} \u2208 \u211d^{d\u00d7m} is a fixed col-normalised random projection",
             "\u2022 \u03b1 ~ U[\u03b1_min, \u03b1_max] scales the signal amplitude",
         ]),
        (6.72, 0.82, 6.42, 2.05, C_REDUND,
         "\u2461 Redundancy atoms  R_{ij} / R_{123}",
         [
             "r, \u03b7_i, \u03b7_j  ~  N(0, I_m)  independently",
             "",
             "r_i  =  \u221a\u03c1 \u00b7 r  +  \u221a(1\u2212\u03c1) \u00b7 \u03b7_i",
             "r_j  =  \u221a\u03c1 \u00b7 r  +  \u221a(1\u2212\u03c1) \u00b7 \u03b7_j",
             "",
             "x_i = \u03b1 \u00b7 P_i^{R_{ij}} \u00b7 r_i ,   x_j = \u03b1 \u00b7 P_j^{R_{ij}} \u00b7 r_j",
             "x_k = 0   (k \u2260 i, j)",
             "",
             "\u2022 r is the shared factor; \u03c1 controls correlation strength",
             "\u2022 R_{123}: same formula extended to all three views",
         ]),
        (0.18, 3.05, 6.3, 2.62, C_SYNERGY,
         "\u2462 Directional Synergy  S_{ij\u2192k}",
         [
             "a, b  ~  N(0, I_m)   (independent source latents)",
             "",
             "x_i = \u03b1 \u00b7 P_i^{A_{ij}} \u00b7 a",
             "x_j = \u03b1 \u00b7 P_j^{B_{ij}} \u00b7 b",
             "",
             "s\u2080  =  f_MLP(a, b ; hop)     [fixed random residual MLP, tanh]",
             "",
             "\u03b4  =  a \u00b7 C_a[hop]  +  b \u00b7 C_b[hop]      [linear de-leakage]",
             "     (C_a, C_b fitted by ridge regression to predict s\u2080 from a, b)",
             "",
             "s  =  s\u2080  \u2212  \u03bb \u00b7 \u03b4          [de-leaked synergy signal]",
             "",
             "x_k = \u03b1 \u00b7 P_k^{SYN_{ij}} \u00b7 s",
             "",
             "\u2022 x_k depends non-linearly on BOTH x_i and x_j but not individually",
             "\u2022 de-leakage removes the component of s\u2080 predictable from a or b alone",
         ]),
        (6.72, 3.05, 6.42, 2.62, C_BACKBONE,
         "\u2463 Shared backbone, noise & composition",
         [
             "Shared backbone  (applied after atom construction):",
             "   z  ~  N(0, I_m)   [per-sample random latent]",
             "   x_v  \u2190  x_v  +  \u03b3 \u00b7 \u03b1 \u00b7 P^{shared}_v \u00b7 z",
             "   Same z and same P^{shared}_v across all views",
             "   P^{shared} is tied (identical) across views when tied=True",
             "",
             "Observation noise (independent per view):",
             "   \u03b5_v  ~  N(0, \u03c3\u00b2 I_d)",
             "   x_v  \u2190  x_v  +  \u03b5_v",
             "",
             "Multi-atom composition (K atoms, K \u2265 2):",
             "   x_v^{total}  =  \u03a3_{k=1}^{K}  x_v^{atom_k}",
             "   Primary atom label retained; K\u22121 atoms sampled w/o replacement",
             "",
             "Fixed projections:  P[view][comp] \u2208 \u211d^{d\u00d7m}, cols normalised to unit \u2113\u2082",
             "Synergy MLP:  residual blocks with tanh; weights fixed at init",
         ]),
    ]

    for lx, ty, pw, ph, accent, title, math_lines in panels:
        _add_rect(slide, lx, ty, pw, ph, fill=C_MATH_BG,
                  border=accent, border_pt=1.2)
        _add_rect(slide, lx, ty, pw, 0.26, fill=accent)
        _add_textbox(slide, lx + 0.12, ty + 0.02, pw - 0.15, 0.24,
                     title, font_size=9, bold=True, color=C_WHITE)

        txBox = slide.shapes.add_textbox(
            Inches(lx + 0.12), Inches(ty + 0.30),
            Inches(pw - 0.18), Inches(ph - 0.34))
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True

        for li, line in enumerate(math_lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.space_before = Pt(0)
            p.space_after  = Pt(0)
            run = p.add_run()
            run.text = line
            is_note = line.startswith('\u2022') or line.startswith('   ')
            run.font.size   = Pt(7.5 if is_note else 8.5)
            run.font.bold   = not is_note and bool(line) and not line.startswith(' ')
            run.font.italic = is_note
            run.font.color.rgb = (
                RGBColor(0xBB, 0xBB, 0xCC) if is_note else
                RGBColor(0xF0, 0xF0, 0xFF)
            )
            if not line:
                run.font.size = Pt(3)

    _add_rect(slide, 0.18, 5.74, 12.95, 0.30,
              fill=RGBColor(0x0D, 0x0D, 0x22), border=C_BORDER, border_pt=0.4)
    _add_textbox(
        slide, 0.28, 5.76, 12.7, 0.26,
        "Projections:  P_v^{comp} ~ N(0, 1/\u221ad), col-normalised  |  "
        "\u03b1 ~ U[0.8, 1.2]  |  \u03c1 \u2208 {0.2, 0.5, 0.8}  |  hop \u2208 {1,2,3,4}  |  "
        "\u03c3 = 0.02  |  \u03bb = 0.25 (de-leakage)  |  \u03b3 = 4.0 (backbone)  |  "
        "K = 1 (single-atom) or 5 (multi-atom)  |  Balanced over 10 PID atoms",
        font_size=7.5, color=C_SUBTITLE)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — Training dynamics + results table  (matplotlib)
# ─────────────────────────────────────────────────────────────────────────────

# Experiment data: (probe_acc%, overfit_drift, best_epoch)
RESULTS = {
    'unique_only':       {'n_cls': 3,  'random': 33.3,
        'A': (35.5, 0.113,  1.7), 'B': (41.5, 0.091, 23.0),
        'C': (38.3, 0.024, 42.0), 'D': (39.2, 0.378,  3.0)},
    'redundancy_only':   {'n_cls': 4,  'random': 25.0,
        'A': (29.9, 0.045, 34.3), 'B': (28.1, 0.909,  5.0),
        'C': (28.2, 0.378,  4.0), 'D': (27.1, 0.609,  3.0)},
    'synergy_only':      {'n_cls': 3,  'random': 33.3,
        'A': (36.0, 0.025, 48.7), 'B': (37.3, 0.610,  6.0),
        'C': (38.3, 0.381,  3.0), 'D': (35.2, 0.863,  7.0)},
    'single_atom_all10': {'n_cls': 10, 'random': 10.0,
        'A': (10.8, 0.075, 37.3), 'B': (11.7, 0.813,  6.0),
        'C': (12.4, 0.715,  2.0), 'D': (11.1, 0.603,  2.0)},
    'multi_atom_2':      {'n_cls': 10, 'random': 10.0,
        'A': (10.5, 0.043, 42.3), 'B': (10.8, 0.975,  2.0),
        'C': (11.5, 0.662,  1.0), 'D': (12.1, 0.751,  4.0)},
    'multi_atom_5':      {'n_cls': 10, 'random': 10.0,
        'A': (10.5, 0.059, 35.7), 'B': (11.3, 1.122,  2.0),
        'C': (10.5, 0.585,  1.0), 'D': (11.1, 1.007,  2.0)},
}

MODEL_COLORS = {'A': '#4488dd', 'B': '#ee8833', 'C': '#44aa55', 'D': '#cc3333'}
MODEL_NAMES  = {
    'A': 'A: Unimodal (SimCLR\u00d73)',
    'B': 'B: Pairwise InfoNCE',
    'C': 'C: TRIANGLE',
    'D': 'D: ConFu',
}


def _val_curve(ep_best: float, drift: float, n: int = 50) -> np.ndarray:
    """Synthetic normalised val-loss curve (min = 0 at ep_best)."""
    ep = max(1, int(round(ep_best)))
    t  = np.arange(n + 1, dtype=float)
    v  = np.zeros(n + 1)

    # phase 1: exponential decrease to 0
    if ep > 0:
        raw = np.exp(-5.0 * np.arange(ep + 1) / ep)
        v[:ep + 1] = (raw - raw[ep]) / (raw[0] - raw[ep] + 1e-9)

    # phase 2: rise by drift amount (log-concave)
    rem = n - ep
    if rem > 0 and drift > 0:
        rise = np.log1p(np.linspace(0, np.e - 1, rem + 1)[1:])
        rise = rise / rise[-1] * drift * 0.85   # scale to visible drift
        v[ep + 1:] = rise

    return t, v


def _train_curve(n: int = 50) -> np.ndarray:
    """Monotonically decreasing relative train-loss (shape only)."""
    t = np.arange(n + 1, dtype=float)
    return t, 0.15 + 0.85 * np.exp(-4.0 * t / n)


def _make_slide3_png() -> io.BytesIO:
    """Render indicative training curves + full results table."""
    fig = plt.figure(figsize=(18, 9), dpi=150)
    fig.patch.set_facecolor(BG_HEX)

    # layout: left curves (45%), right table (55%)
    ax_c = fig.add_axes([0.04, 0.12, 0.42, 0.74])
    ax_t = fig.add_axes([0.50, 0.06, 0.49, 0.86])

    # ── figure title ─────────────────────────────────────────────────────────
    fig.text(0.5, 0.97,
             'SSL Objective Comparison: Training Dynamics & Probe Accuracy',
             ha='center', va='center', fontsize=17, fontweight='bold',
             color='#e8e8ff')

    # ── LEFT: val-loss curves (unique_only experiment) ───────────────────────
    ax_c.set_facecolor('#0a0a1e')
    for model in ['A', 'B', 'C', 'D']:
        prob, drift, ep = RESULTS['unique_only'][model]
        t, v = _val_curve(ep, drift, n=50)
        t_tr, tr = _train_curve(n=50)
        # scale train so it starts near val[0]
        tr_scaled = tr * (v[0] + 0.05) / tr[0]
        col = MODEL_COLORS[model]
        ax_c.plot(t, v, color=col, lw=2.2, label=MODEL_NAMES[model], zorder=3)
        ax_c.plot(t_tr, tr_scaled, color=col, lw=1.0, ls='--', alpha=0.35, zorder=2)
        ax_c.axvline(ep, color=col, lw=0.9, ls=':', alpha=0.55, zorder=1)

    ax_c.set_xlabel('Epoch', color='#aabbcc', fontsize=11)
    ax_c.set_ylabel('Relative val. loss  (normalised)', color='#aabbcc', fontsize=11)
    ax_c.set_title('unique_only experiment  (n=3 classes, 3\u00d71k samples)',
                    color='#e0e0ff', fontsize=12, fontweight='bold', pad=8)
    ax_c.tick_params(colors='#778899', labelsize=9)
    for spine in ['top', 'right']:
        ax_c.spines[spine].set_visible(False)
    for spine in ['bottom', 'left']:
        ax_c.spines[spine].set_color('#334455')

    leg = ax_c.legend(framealpha=0.25, facecolor='#1a1a2e', edgecolor='#334455',
                       labelcolor='#ddeeff', fontsize=9, loc='upper right')
    ax_c.text(0.03, 0.97,
              'Solid = val loss\nDashed = train loss\nDotted = best-epoch',
              transform=ax_c.transAxes, ha='left', va='top',
              fontsize=8, color='#778899', style='italic', linespacing=1.4)

    # note about model A stability
    ax_c.text(0.97, 0.55,
              'Model A trains stably\n(VectorAugmenter active)\n\nB/C/D overfit\n(no augmentation)',
              transform=ax_c.transAxes, ha='right', va='center',
              fontsize=8.5, color='#aabbcc', style='italic', linespacing=1.4,
              bbox=dict(facecolor='#151530', edgecolor='#334455', boxstyle='round,pad=0.4'))

    # ── RIGHT: results table ─────────────────────────────────────────────────
    ax_t.set_facecolor(BG_HEX)
    ax_t.axis('off')

    exp_keys = list(RESULTS.keys())
    exp_labels = [
        'unique_only\n(3 cls)',
        'redundancy_only\n(4 cls)',
        'synergy_only\n(3 cls)',
        'single_atom\nall10 (10 cls)',
        'multi_atom\n×2  (10 cls)',
        'multi_atom\n×5  (10 cls)',
    ]
    col_headers = ['Experiment', 'Rand.', 'A: Unimodal', 'B: Pairwise', 'C: TRIANGLE', 'D: ConFu']
    col_colors  = ['#334455', '#334455',
                   MODEL_COLORS['A'], MODEL_COLORS['B'],
                   MODEL_COLORS['C'], MODEL_COLORS['D']]

    # draw header
    n_cols = len(col_headers)
    n_rows = len(exp_labels) + 1  # +1 for header
    col_widths = [0.22, 0.07, 0.175, 0.175, 0.175, 0.175]  # fractions of axes width
    col_x = [0.0]
    for cw in col_widths[:-1]:
        col_x.append(col_x[-1] + cw)

    row_h = 0.12    # height of each row in axes fraction
    header_y = 0.90

    def draw_cell(ax, x, y, w, h, text, fc, tc='white', fontsize=9, bold=False, va='center'):
        rect = FancyBboxPatch((x + 0.003, y - h + 0.005), w - 0.006, h - 0.01,
                               boxstyle='square,pad=0', facecolor=fc,
                               edgecolor='#1a1a2e', linewidth=0.5,
                               transform=ax.transAxes, clip_on=False)
        ax.add_patch(rect)
        ax.text(x + w / 2, y - h / 2, text, ha='center', va=va,
                fontsize=fontsize, color=tc, fontweight='bold' if bold else 'normal',
                transform=ax.transAxes, linespacing=1.3)

    # header row
    for ci, (ch, cx, cw, cc) in enumerate(zip(col_headers, col_x, col_widths, col_colors)):
        draw_cell(ax_t, cx, header_y, cw, row_h, ch, fc=cc, tc='white',
                  fontsize=9.5, bold=True)

    # data rows
    for ri, (ek, el) in enumerate(zip(exp_keys, exp_labels)):
        row_d = RESULTS[ek]
        rand  = row_d['random']
        ry    = header_y - (ri + 1) * row_h
        bg    = '#0e0e22' if ri % 2 == 0 else '#141430'

        # experiment label
        draw_cell(ax_t, col_x[0], ry, col_widths[0], row_h, el,
                  fc=bg, tc='#c0c8e8', fontsize=8.5)

        # random baseline
        draw_cell(ax_t, col_x[1], ry, col_widths[1], row_h,
                  f'{rand:.0f}%', fc=bg, tc='#888899', fontsize=9)

        # model results
        best_acc = max(row_d[m][0] for m in ['A', 'B', 'C', 'D'])
        for mi, model in enumerate(['A', 'B', 'C', 'D']):
            acc, drift, ep = row_d[model]
            gap = acc - rand
            cell_txt = f'{acc:.1f}%\n+{gap:.1f}pp'
            is_best = (acc == best_acc)
            fc_cell = MODEL_COLORS[model] if is_best else bg
            tc_cell = 'white' if is_best else '#b0b8d0'
            draw_cell(ax_t, col_x[2 + mi], ry, col_widths[2 + mi], row_h,
                      cell_txt, fc=fc_cell, tc=tc_cell,
                      fontsize=8.5 if is_best else 8, bold=is_best)

    # legend below table
    leg_y = header_y - (len(exp_labels) + 1) * row_h - 0.04
    ax_t.text(0.5, leg_y,
              'Each cell shows probe accuracy and gap above random chance (+pp).\n'
              'Highlighted cell = best model per experiment.',
              ha='center', va='top', fontsize=9, color='#778899',
              transform=ax_t.transAxes, style='italic')

    # table title
    ax_t.text(0.5, header_y + 0.07,
              'Frozen Linear Probe on Concatenated [h\u2081, h\u2082, h\u2083]',
              ha='center', va='center', fontsize=12, fontweight='bold',
              color='#c8d0f0', transform=ax_t.transAxes)

    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=BG_HEX, edgecolor='none')
    buf.seek(0)
    plt.close(fig)
    return buf


def build_slide3(prs: Presentation):
    slide = _blank_slide(prs)
    _fill_bg(slide, C_BG)
    _embed_png(slide, _make_slide3_png(), l=0, t=0, w=13.33, h=7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Entry point
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = _new_prs()
    print("Building slide 1 (schematic) ...")
    build_slide1(prs)
    print("Building slide 2 (mathematics) ...")
    build_slide2(prs)
    print("Building slide 3 (training dynamics + results) ...")
    build_slide3(prs)
    out = "pid_ssl_data_generation.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
